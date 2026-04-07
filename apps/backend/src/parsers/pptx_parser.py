"""PPTX document parser."""
import os
from typing import List

from pptx import Presentation
from pptx.util import Pt

from src.core.logging import get_logger
from src.graphs.state import DocumentMetadata, ExtractedSection
from src.parsers.base_parser import BaseParser, ParseResult

logger = get_logger(__name__)


class PPTXParser(BaseParser):
    """Parser for PowerPoint (.pptx) files."""

    def can_parse(self, file_path: str) -> bool:
        return file_path.lower().endswith((".pptx", ".ppt"))

    def parse(self, file_path: str) -> ParseResult:
        logger.info("parsing_pptx", file_path=file_path)
        warnings: List[str] = []
        sections: List[ExtractedSection] = []
        content_parts: List[str] = []

        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts: List[str] = []
                title_text = ""

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # Detect title by placeholder type or large font
                        is_title = (
                            hasattr(shape, "placeholder_format")
                            and shape.placeholder_format is not None
                            and shape.placeholder_format.idx in (0, 1)
                        )
                        if is_title and not title_text:
                            title_text = text
                        else:
                            slide_texts.append(text)

                slide_content = "\n".join(slide_texts)
                if slide_content or title_text:
                    label = title_text or f"Slide {slide_num}"
                    content_parts.append(f"## {label}\n{slide_content}")
                    sections.append(ExtractedSection(
                        title=label,
                        content=slide_content,
                        level=1,
                        page=slide_num,
                    ))

            full_content = "\n\n".join(content_parts)
            normalized = self.normalize_text(full_content)

            metadata = DocumentMetadata(
                filename=os.path.basename(file_path),
                file_type="pptx",
                file_size=os.path.getsize(file_path),
                page_count=len(prs.slides),
                word_count=len(normalized.split()),
            )

            return ParseResult(
                content=normalized,
                metadata=metadata,
                sections=sections,
                tables=[],
                warnings=warnings,
            )

        except Exception as e:
            logger.error("pptx_parsing_failed", error=str(e))
            raise ValueError(f"Failed to parse PPTX: {str(e)}")
